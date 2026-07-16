import os
import logging
from typing import Optional, Dict, Any
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from typing import Dict, Any, List

logger = logging.getLogger("a3.agents.reporting")

class ReportingAgent:
    @classmethod
    def generate_pdf_report(
        cls, 
        dataset_name: str, 
        metadata: Dict[str, Any], 
        quality_report: Dict[str, Any],
        insights: Dict[str, Any], 
        forecast_commentary: Optional[str],
        output_path: str
    ) -> str:
        """Generates a highly structured executive PDF report using ReportLab."""
        try:
            # Setup document
            doc = SimpleDocTemplate(
                output_path, 
                pagesize=letter,
                rightMargin=54, leftMargin=54,
                topMargin=54, bottomMargin=54
            )
            
            story = []
            
            # Setup styles
            styles = getSampleStyleSheet()
            
            # Custom styled paragraph definitions
            title_style = ParagraphStyle(
                'DocTitle',
                parent=styles['Heading1'],
                fontName='Helvetica-Bold',
                fontSize=26,
                leading=32,
                textColor=colors.HexColor('#1e1b4b'), # Deep Indigo
                spaceAfter=15
            )
            
            subtitle_style = ParagraphStyle(
                'DocSubtitle',
                parent=styles['Normal'],
                fontName='Helvetica-Oblique',
                fontSize=12,
                leading=16,
                textColor=colors.HexColor('#475569'), # Slate grey
                spaceAfter=30
            )

            h1_style = ParagraphStyle(
                'SectionHeader',
                parent=styles['Heading2'],
                fontName='Helvetica-Bold',
                fontSize=16,
                leading=20,
                textColor=colors.HexColor('#312e81'),
                spaceBefore=15,
                spaceAfter=10,
                keepWithNext=True
            )

            h2_style = ParagraphStyle(
                'SubSectionHeader',
                parent=styles['Heading3'],
                fontName='Helvetica-Bold',
                fontSize=12,
                leading=16,
                textColor=colors.HexColor('#4f46e5'),
                spaceBefore=8,
                spaceAfter=4,
                keepWithNext=True
            )

            body_style = ParagraphStyle(
                'BodyTextCustom',
                parent=styles['Normal'],
                fontName='Helvetica',
                fontSize=10,
                leading=14,
                textColor=colors.HexColor('#1e293b'),
                spaceAfter=8
            )

            bullet_style = ParagraphStyle(
                'BulletCustom',
                parent=styles['Normal'],
                fontName='Helvetica',
                fontSize=10,
                leading=14,
                textColor=colors.HexColor('#1e293b'),
                leftIndent=15,
                firstLineIndent=-10,
                spaceAfter=6
            )
            
            # --- COVER PAGE CONTENT ---
            story.append(Spacer(1, 1 * inch))
            story.append(Paragraph("A3 Analytics Platform", subtitle_style))
            story.append(Paragraph("EXECUTIVE DATA INTELLIGENCE REPORT", title_style))
            story.append(Paragraph(f"Dataset: {dataset_name}", ParagraphStyle('MetaName', parent=body_style, fontSize=11, fontName='Helvetica-Bold')))
            
            import datetime
            date_str = datetime.datetime.now().strftime("%B %d, %Y")
            story.append(Paragraph(f"Generated on: {date_str}", body_style))
            story.append(Paragraph("Classification: Private & Confidential", ParagraphStyle('Classify', parent=body_style, textColor=colors.HexColor('#ef4444'), fontName='Helvetica-Bold')))
            story.append(Spacer(1, 0.5 * inch))
            
            # Simple decorative divider
            divider = Table([['']], colWidths=[504])
            divider.setStyle(TableStyle([
                ('LINEBELOW', (0,0), (-1,-1), 3, colors.HexColor('#4f46e5')),
                ('BOTTOMPADDING', (0,0), (-1,-1), 0),
                ('TOPPADDING', (0,0), (-1,-1), 0)
            ]))
            story.append(divider)
            story.append(Spacer(1, 0.4 * inch))
            
            story.append(Paragraph("<b>Executive Summary</b>", h1_style))
            summary_text = metadata.get("summary", "This report contains structural and analytical insight from the dataset.")
            story.append(Paragraph(summary_text, body_style))
            story.append(PageBreak())
            
            # --- PAGE 2: METADATA & DATA QUALITY ---
            story.append(Paragraph("1. Dataset Profile & Quality Scorecard", h1_style))
            
            # Create a nice metadata table
            meta_data_table = [
                [Paragraph("<b>Attribute</b>", body_style), Paragraph("<b>Value</b>", body_style)],
                ["Business Domain", metadata.get("business_domain", "Unknown")],
                ["File Type", metadata.get("file_type", "Unknown").upper()],
                ["Total Row Count", f"{metadata.get('row_count', 0):,}"],
                ["Total Column Count", str(metadata.get('column_count', 0))],
                ["File Size", f"{metadata.get('file_size', 0) / 1024:.2f} KB"]
            ]
            t_meta = Table(meta_data_table, colWidths=[200, 304])
            t_meta.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#f1f5f9')),
                ('BOTTOMPADDING', (0,0), (-1,-1), 6),
                ('TOPPADDING', (0,0), (-1,-1), 6),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#cbd5e1')),
            ]))
            story.append(t_meta)
            story.append(Spacer(1, 15))
            
            # Quality Score Card
            q_score = quality_report.get("score", 100.0)
            score_color = '#10b981' if q_score >= 85 else ('#f59e0b' if q_score >= 60 else '#ef4444')
            story.append(Paragraph(f"Data Quality Score: <font color='{score_color}'><b>{q_score}/100</b></font>", h2_style))
            story.append(Paragraph(quality_report.get("narrative_summary", ""), body_style))
            story.append(Spacer(1, 10))
            
            # Quality Issues Table
            issues = quality_report.get("issues", [])
            if issues:
                story.append(Paragraph("Identified Data Issues", h2_style))
                issue_rows = [[Paragraph("<b>Severity</b>", body_style), Paragraph("<b>Issue Description</b>", body_style)]]
                for issue in issues[:8]: # Cap list
                    sev = issue.get("severity", "low").upper()
                    sev_color = '#ef4444' if sev == 'HIGH' else ('#f59e0b' if sev == 'MEDIUM' else '#3b82f6')
                    issue_rows.append([
                        Paragraph(f"<font color='{sev_color}'><b>{sev}</b></font>", body_style),
                        Paragraph(issue.get("message", ""), body_style)
                    ])
                t_issues = Table(issue_rows, colWidths=[100, 404])
                t_issues.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#f8fafc')),
                    ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
                    ('BOTTOMPADDING', (0,0), (-1,-1), 5),
                    ('TOPPADDING', (0,0), (-1,-1), 5),
                ]))
                story.append(t_issues)
            else:
                story.append(Paragraph("✓ No major data quality issues were found. The dataset is fully consistent.", body_style))
                
            story.append(PageBreak())
            
            # --- PAGE 3: INSIGHTS & OPPORTUNITIES ---
            story.append(Paragraph("2. AI-Generated Insights & Recommendations", h1_style))
            
            story.append(Paragraph("Key Findings", h2_style))
            for finding in insights.get("key_findings", []):
                story.append(Paragraph(f"• {finding}", bullet_style))
                
            story.append(Spacer(1, 10))
            story.append(Paragraph("Strategic Business Opportunities", h2_style))
            for opp in insights.get("business_opportunities", []):
                story.append(Paragraph(f"• <b>Opportunity:</b> {opp}", bullet_style))
                
            story.append(Spacer(1, 10))
            story.append(Paragraph("Key Risks & Anomaly Indicators", h2_style))
            for risk in insights.get("risks", []):
                story.append(Paragraph(f"• <b>Risk Alert:</b> {risk}", bullet_style))
                
            if forecast_commentary:
                story.append(Spacer(1, 15))
                story.append(Paragraph("3. Predictive Forecast & Growth Summary", h1_style))
                story.append(Paragraph(forecast_commentary, body_style))
                
            # Build Document
            doc.build(story)
            return output_path
            
        except Exception as e:
            logger.error(f"Error compiling PDF report: {e}")
            raise e

    @classmethod
    def generate_pptx_report(
        cls, 
        dataset_name: str, 
        presentation_deck: List[Dict[str, Any]], 
        output_path: str
    ) -> str:
        """Generates a professional slide deck presentation using python-pptx."""
        try:
            from pptx import Presentation
            
            prs = Presentation()
            
            # Slide 1: Title Slide
            slide_layout = prs.slide_layouts[0] # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = f"Executive Presentation Briefing"
            subtitle.text = f"Dataset: {dataset_name}\nCreated autonomously by A3 platform"
            
            # Slides from presentation_deck
            for s in presentation_deck:
                slide_layout = prs.slide_layouts[1] # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Title
                slide.shapes.title.text = s.get("title", "Insight Slide")
                
                # Content placeholder
                content_placeholder = slide.placeholders[1]
                
                text_parts = []
                if s.get("subtitle"):
                    text_parts.append(f"Focus Area: {s['subtitle']}\n")
                    
                bullets = s.get("bullets", [])
                if bullets:
                    for b in bullets:
                        text_parts.append(f"• {b}")
                else:
                    metrics = s.get("metrics", [])
                    for m in metrics:
                        text_parts.append(f"• {m.get('label')}: {m.get('value')}")
                        
                content_placeholder.text = "\n".join(text_parts)
                
            prs.save(output_path)
            return output_path
        except Exception as e:
            logger.error(f"Error compiling PPTX report: {e}")
            raise e
